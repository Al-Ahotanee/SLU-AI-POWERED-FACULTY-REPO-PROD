"""
╔══════════════════════════════════════════════════════════════════════╗
║   AI-Powered Faculty Repository System                              ║
║   Kafin Hausa State University · Dept. of Computer Science         ║
║   Stack: Flask + SQLite + React 18 + Tailwind CSS + Lucide Icons   ║
║                                                                     ║
║   AI ENGINE: Hugging Face Space (sentence-transformers)            ║
║   Set env var  HF_SPACE_URL = https://<user>-<space>.hf.space      ║
║   Optional:    HF_API_KEY   = <your-secret-key>                    ║
╚══════════════════════════════════════════════════════════════════════╝
"""

import os, json, re, math, uuid, logging, threading
from datetime import datetime, timedelta, timezone
def _utcnow(): return datetime.now(timezone.utc)
from functools import wraps
from pathlib import Path
from collections import defaultdict

from flask import Flask, request, jsonify, send_from_directory, g
from flask_cors import CORS
import sqlite3, jwt, bcrypt

# ── HTTP client for HF Space ───────────────────────────────────────────────────
import requests as _req
_HF_URL     = os.environ.get("HF_SPACE_URL", "").rstrip("/")
_HF_KEY     = os.environ.get("HF_API_KEY", "")
_HF_TIMEOUT = int(os.environ.get("HF_TIMEOUT", 10))   # seconds

def _hf(endpoint: str, payload: dict):
    """POST to HF Space and return the JSON response dict, or None on failure."""
    if not _HF_URL:
        return None
    url = f"{_HF_URL}{endpoint}"
    headers = {"Content-Type": "application/json"}
    if _HF_KEY:
        headers["X-API-Key"] = _HF_KEY
    try:
        r = _req.post(url, json=payload, headers=headers, timeout=_HF_TIMEOUT)
        r.raise_for_status()
        return r.json()
    except Exception as e:
        log.warning(f"HF Space call to {endpoint} failed: {e}")
        return None

# ── optional document-parsing libs ────────────────────────────────────────────
try:
    import fitz;                    HAS_PDF  = True
except ImportError:                 HAS_PDF  = False
try:
    from docx import Document as _DocxDoc; HAS_DOCX = True
except ImportError:                 HAS_DOCX = False
try:
    from pptx import Presentation;  HAS_PPTX = True
except ImportError:                 HAS_PPTX = False

# ── config ────────────────────────────────────────────────────────────────────
BASE_DIR   = Path(__file__).parent
UPLOAD_DIR = Path(os.environ.get("UPLOAD_DIR", str(BASE_DIR / "uploads")))
DB_PATH    = Path(os.environ.get("DB_PATH",    str(BASE_DIR / "repository.db")))
SECRET     = os.environ.get("SECRET_KEY", "slu-faculty-repo-secret-key-2024!!")
JWT_DAYS   = int(os.environ.get("JWT_DAYS", 7))
MAX_BYTES  = int(os.environ.get("MAX_MB", 50)) * 1024 * 1024
PORT       = int(os.environ.get("PORT", 5000))
ALLOWED    = {".pdf",".docx",".pptx",".txt",".xlsx",".png",".jpg",".jpeg"}

UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
logging.basicConfig(level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s")
log = logging.getLogger("SLURepo")

app = Flask(__name__, static_folder=str(BASE_DIR), static_url_path="")
app.config["MAX_CONTENT_LENGTH"] = MAX_BYTES
CORS(app, resources={r"/api/*": {"origins": "*"}},
     supports_credentials=True,
     allow_headers=["Content-Type","Authorization"])
init_db()
# ── db ────────────────────────────────────────────────────────────────────────
def get_db():
    if "db" not in g:
        g.db = sqlite3.connect(str(DB_PATH), detect_types=sqlite3.PARSE_DECLTYPES,
                               check_same_thread=False)
        g.db.row_factory = sqlite3.Row
        g.db.execute("PRAGMA foreign_keys = ON")
        g.db.execute("PRAGMA journal_mode = WAL")
    return g.db

@app.teardown_appcontext
def close_db(e=None):
    db = g.pop("db", None)
    if db: db.close()

def init_db():
    with app.app_context():
        db = get_db()
        db.executescript("""
        CREATE TABLE IF NOT EXISTS users (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT NOT NULL, email TEXT UNIQUE NOT NULL,
            password TEXT NOT NULL, role TEXT DEFAULT 'faculty',
            department TEXT DEFAULT '', bio TEXT DEFAULT '',
            active INTEGER DEFAULT 1,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            last_login TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS documents (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            title TEXT NOT NULL, description TEXT DEFAULT '',
            filename TEXT NOT NULL, original_name TEXT NOT NULL,
            file_type TEXT NOT NULL, file_size INTEGER DEFAULT 0,
            course_code TEXT DEFAULT '', academic_level TEXT DEFAULT '',
            resource_type TEXT DEFAULT '', academic_year TEXT DEFAULT '',
            uploader_id INTEGER NOT NULL, is_public INTEGER DEFAULT 1,
            text_content TEXT DEFAULT '', embedding TEXT DEFAULT '',
            download_count INTEGER DEFAULT 0, view_count INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (uploader_id) REFERENCES users(id)
        );
        CREATE TABLE IF NOT EXISTS tags (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL, tag TEXT NOT NULL,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS bookmarks (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER NOT NULL, document_id INTEGER NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            UNIQUE(user_id, document_id),
            FOREIGN KEY (user_id) REFERENCES users(id) ON DELETE CASCADE,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE
        );
        CREATE TABLE IF NOT EXISTS comments (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document_id INTEGER NOT NULL, user_id INTEGER NOT NULL,
            content TEXT NOT NULL,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (document_id) REFERENCES documents(id) ON DELETE CASCADE,
            FOREIGN KEY (user_id) REFERENCES users(id)
        );
        CREATE TABLE IF NOT EXISTS search_logs (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER, query TEXT NOT NULL, result_count INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        CREATE TABLE IF NOT EXISTS activity_log (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id INTEGER, action TEXT NOT NULL, document_id INTEGER,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        CREATE INDEX IF NOT EXISTS idx_d_uploader ON documents(uploader_id);
        CREATE INDEX IF NOT EXISTS idx_d_course   ON documents(course_code);
        CREATE INDEX IF NOT EXISTS idx_t_doc      ON tags(document_id);
        CREATE INDEX IF NOT EXISTS idx_t_tag      ON tags(tag);
        CREATE INDEX IF NOT EXISTS idx_bm         ON bookmarks(user_id);
        CREATE INDEX IF NOT EXISTS idx_cmt        ON comments(document_id);
        """)
        db.commit()
        if not db.execute("SELECT id FROM users WHERE role='admin' LIMIT 1").fetchone():
            pw = bcrypt.hashpw(b"Admin@1234", bcrypt.gensalt()).decode()
            db.execute("INSERT INTO users (name,email,password,role,department) VALUES (?,?,?,?,?)",
                       ("System Admin","admin@slu.edu.ng", pw,"admin","Computer Science"))
            db.commit()
            log.info("Default admin: admin@slu.edu.ng / Admin@1234")

# ── auth helpers ──────────────────────────────────────────────────────────────
def make_token(uid, role):
    t = jwt.encode({"sub":uid,"role":role,
                    "exp": _utcnow()+timedelta(days=JWT_DAYS),
                    "iat": _utcnow()}, SECRET, algorithm="HS256")
    return t if isinstance(t, str) else t.decode()

def require_auth(f):
    @wraps(f)
    def w(*a, **kw):
        auth = request.headers.get("Authorization","")
        if not auth.startswith("Bearer "):
            return jsonify({"error":"Missing token"}), 401
        try:
            d = jwt.decode(auth.split(" ",1)[1], SECRET, algorithms=["HS256"])
            g.user_id = d["sub"]; g.user_role = d["role"]
        except jwt.ExpiredSignatureError:
            return jsonify({"error":"Token expired"}), 401
        except Exception:
            return jsonify({"error":"Invalid token"}), 401
        return f(*a, **kw)
    return w

def require_admin(f):
    @wraps(f)
    @require_auth
    def w(*a, **kw):
        if g.user_role != "admin":
            return jsonify({"error":"Admin only"}), 403
        return f(*a, **kw)
    return w

# ── text extraction ───────────────────────────────────────────────────────────
def extract_text(path, ext):
    try:
        if ext == ".pdf"  and HAS_PDF:
            return " ".join(p.get_text() for p in fitz.open(str(path)))[:50000]
        if ext == ".docx" and HAS_DOCX:
            return " ".join(p.text for p in _DocxDoc(str(path)).paragraphs)[:50000]
        if ext == ".pptx" and HAS_PPTX:
            prs = Presentation(str(path))
            return " ".join(s.text for sl in prs.slides
                            for s in sl.shapes if hasattr(s,"text"))[:50000]
        if ext == ".txt":
            return path.read_text(errors="ignore")[:50000]
    except Exception as e:
        log.warning(f"extract_text: {e}")
    return ""

# ── keyword NLP (local, zero-dep fallback) ────────────────────────────────────
COURSE_KW = {
    "CSC101":["programming","algorithm","introduction","computer","basic"],
    "CSC201":["data structure","stack","queue","linked list","tree","binary"],
    "CSC301":["database","sql","relational","normalization","query"],
    "CSC302":["operating system","process","memory","scheduling","thread"],
    "CSC401":["artificial intelligence","machine learning","neural","deep learning"],
    "CSC402":["network","protocol","tcp","ip","routing","socket"],
    "CSC403":["software engineering","sdlc","agile","design pattern","uml"],
    "CSC404":["computer graphics","rendering","opengl","pixel","3d"],
    "CSC501":["research","methodology","thesis","dissertation"],
    "MAT101":["mathematics","calculus","algebra","differential","integral"],
    "STA101":["statistics","probability","regression","hypothesis"],
}
RES_KW = {
    "Lecture Note":  ["lecture","note","introduction","overview","chapter"],
    "Research Paper":["abstract","methodology","conclusion","references","journal"],
    "Assignment":    ["assignment","exercise","problem","submit","question"],
    "Past Question": ["past question","examination","exam","test","quiz"],
    "Textbook":      ["textbook","edition","publisher","isbn"],
    "Presentation":  ["slide","presentation","powerpoint"],
    "Tutorial":      ["tutorial","guide","step by step","how to"],
    "Project":       ["project","final year","implementation","system design"],
}
LVL_KW = {
    "100 Level":  ["100 level","year one","freshman","first year"],
    "200 Level":  ["200 level","year two","sophomore"],
    "300 Level":  ["300 level","year three","junior"],
    "400 Level":  ["400 level","year four","senior","final year"],
    "Postgraduate":["postgraduate","masters","phd","msc"],
}
STOPS = {"the","a","an","and","or","but","in","on","at","to","for","of","with",
         "is","are","was","were","be","been","have","has","had","do","does",
         "did","will","would","could","should","that","this","those","it","its",
         "we","our","they","their","which","from","by","as","not","no","all"}

def _kw_score(t, km):
    s = {l: sum(t.count(k) for k in ks) for l,ks in km.items()}
    s = {l:v for l,v in s.items() if v>0}
    return max(s,key=s.get) if s else None

def auto_classify_local(text, fname):
    """Pure-keyword fallback classifier (no network needed)."""
    t = (text + " " + fname).lower()
    return {"course_code":  _kw_score(t,COURSE_KW) or "",
            "resource_type":_kw_score(t,RES_KW) or "General",
            "academic_level":_kw_score(t,LVL_KW) or ""}

def auto_classify(text, fname):
    """Try HF Space first; fall back to local keyword classifier."""
    res = _hf("/classify", {"text": text[:3000], "filename": fname})
    if res:
        return res
    return auto_classify_local(text, fname)

def extract_tags(text, n=12):
    words = re.findall(r'\b[a-zA-Z]{4,}\b', text.lower())
    freq  = defaultdict(int)
    for w in words:
        if w not in STOPS: freq[w] += 1
    for i in range(len(words)-1):
        a,b = words[i],words[i+1]
        if a not in STOPS and b not in STOPS: freq[f"{a} {b}"] += 2
    return [w for w,_ in sorted(freq.items(),key=lambda x:x[1],reverse=True)[:n]]

def get_embedding(text):
    """Encode text via HF Space; returns list[float] or None."""
    if not text.strip():
        return None
    res = _hf("/encode", {"texts": [text[:5000]]})
    if res and res.get("embeddings"):
        return res["embeddings"][0]
    return None

# ── TF-IDF fallback (no external deps) ───────────────────────────────────────
def tfidf(corpus):
    N   = len(corpus)
    tok = [re.findall(r'\b\w+\b', d.lower()) for d in corpus]
    df  = defaultdict(int)
    for ts in tok:
        for t in set(ts): df[t] += 1
    vecs = []
    for ts in tok:
        tf = defaultdict(int)
        for t in ts: tf[t] += 1
        vecs.append({t:(c/(len(ts) or 1))*(math.log((N+1)/(df[t]+1))+1)
                     for t,c in tf.items()})
    return vecs

def cos_dict(v1,v2):
    keys = set(v1)&set(v2)
    if not keys: return 0.0
    dot = sum(v1[k]*v2[k] for k in keys)
    return dot/(math.sqrt(sum(x*x for x in v1.values()))*
                math.sqrt(sum(x*x for x in v2.values()))+1e-9)

# ── search ────────────────────────────────────────────────────────────────────
def search_docs(query, filters, uid=None):
    db = get_db()
    db.execute("INSERT INTO search_logs (user_id,query) VALUES (?,?)",(uid,query)); db.commit()
    sql = ("SELECT d.*,u.name as uploader_name,GROUP_CONCAT(t.tag,',') as tags "
           "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
           "LEFT JOIN tags t ON t.document_id=d.id WHERE d.is_public=1")
    p = []
    if filters.get("course"):  sql+=" AND d.course_code=?";   p.append(filters["course"])
    if filters.get("year"):    sql+=" AND d.academic_year=?"; p.append(filters["year"])
    if filters.get("type"):    sql+=" AND d.resource_type=?"; p.append(filters["type"])
    if filters.get("level"):   sql+=" AND d.academic_level=?";p.append(filters["level"])
    if filters.get("author") and str(filters["author"]).isdigit():
        sql+=" AND d.uploader_id=?"; p.append(int(filters["author"]))
    sql += " GROUP BY d.id"
    docs = [dict(r) for r in db.execute(sql,p).fetchall()]

    if not query.strip():
        docs.sort(key=lambda x:x["created_at"],reverse=True); return docs[:50]

    ql = query.lower()

    # ── Try semantic search via HF Space ─────────────────────────────────────
    # Build list of candidates that have stored embeddings
    candidates_with_emb = []
    for d in docs:
        if d.get("embedding"):
            try:
                emb = json.loads(d["embedding"])
                candidates_with_emb.append({"id": d["id"], "embedding": emb})
            except Exception:
                pass

    hf_scores = {}
    if candidates_with_emb and _HF_URL:
        res = _hf("/search", {"query": query, "candidates": candidates_with_emb})
        if res and res.get("results"):
            for item in res["results"]:
                hf_scores[item["id"]] = item["score"]

    if hf_scores:
        scored = []
        for d in docs:
            kw = 0
            for f in ["title","description","course_code","tags"]:
                val = str(d.get(f,"") or "").lower()
                if ql in val: kw += 3
                for w in ql.split():
                    if w in val: kw += 1
            sem = hf_scores.get(d["id"], 0.0) * 10
            tot = sem + kw
            if tot > 0:
                d["relevance"] = round(tot, 4)
                scored.append(d)
        return sorted(scored, key=lambda x: x["relevance"], reverse=True)[:50]

    # ── TF-IDF fallback ───────────────────────────────────────────────────────
    corpus = [f"{d['title']} {d['description']} {d['text_content']}" for d in docs]
    if not corpus: return []
    vecs = tfidf(corpus + [query]); qv = vecs[-1]; scored = []
    for i, d in enumerate(docs):
        sim = cos_dict(vecs[i], qv)
        if ql in str(d.get("title","")).lower(): sim += 0.5
        if sim > 0.01: d["relevance"] = round(sim,4); scored.append(d)
    return sorted(scored, key=lambda x: x["relevance"], reverse=True)[:50]


def get_recs(doc_id, n=6):
    db = get_db()
    target = db.execute("SELECT * FROM documents WHERE id=?",(doc_id,)).fetchone()
    if not target: return []
    target = dict(target)
    cands  = [dict(r) for r in db.execute(
        "SELECT d.*,GROUP_CONCAT(t.tag,',') as tags FROM documents d "
        "LEFT JOIN tags t ON t.document_id=d.id WHERE d.id!=? AND d.is_public=1 GROUP BY d.id",
        (doc_id,)).fetchall()]
    if not cands: return []

    # ── Try HF Space recommendations ──────────────────────────────────────────
    if target.get("embedding") and _HF_URL:
        try:
            te = json.loads(target["embedding"])
            hf_cands = []
            for c in cands:
                item = {"id": c["id"],
                        "course_code": c.get("course_code",""),
                        "academic_level": c.get("academic_level","")}
                if c.get("embedding"):
                    try: item["embedding"] = json.loads(c["embedding"])
                    except: pass
                hf_cands.append(item)
            res = _hf("/recommend", {
                "target_embedding": te,
                "target_course":    target.get("course_code",""),
                "target_level":     target.get("academic_level",""),
                "candidates":       hf_cands
            })
            if res and res.get("results"):
                score_map = {item["id"]: item["score"] for item in res["results"]}
                results = []
                for c in cands:
                    s = score_map.get(c["id"], 0.0)
                    results.append((c, s))
                results.sort(key=lambda x: x[1], reverse=True)
                return [dict(d, relevance=round(s,4)) for d,s in results[:n]]
        except Exception as e:
            log.warning(f"HF recommend failed: {e}")

    # ── TF-IDF fallback ───────────────────────────────────────────────────────
    corpus = [f"{d['title']} {d['description']} {d['text_content']}" for d in cands]
    tt     = f"{target['title']} {target['description']} {target['text_content']}"
    vecs   = tfidf(corpus + [tt]); tv = vecs[-1]
    scores = []
    for i, c in enumerate(cands):
        sim = cos_dict(vecs[i], tv)
        if c["course_code"] == target["course_code"] and c["course_code"]: sim += 0.3
        scores.append((c, sim))
    scores.sort(key=lambda x: x[1], reverse=True)
    return [dict(d, relevance=round(s,4)) for d,s in scores[:n]]


# ═══════════════════════════════════════════════════════════════════════════════
#  ROUTES
# ═══════════════════════════════════════════════════════════════════════════════

# SPA catch-all
@app.route("/",defaults={"path":""})
@app.route("/<path:path>")
def spa(path):
    if path.startswith("api/") or path.startswith("uploads/"): return jsonify({"error":"Not found"}),404
    return send_from_directory(str(BASE_DIR),"index.html")

# Health
@app.route("/api/health")
def health():
    db = get_db()
    # Check HF Space status
    hf_status = "not configured"
    if _HF_URL:
        try:
            r = _req.get(f"{_HF_URL}/health", timeout=5)
            hf_status = "online" if r.ok else "error"
        except Exception:
            hf_status = "offline"
    ai_mode = "semantic (HF Space)" if hf_status == "online" else "tfidf (fallback)"
    return jsonify({
        "status":    "ok",
        "ai":        ai_mode,
        "hf_space":  hf_status,
        "documents": db.execute("SELECT COUNT(*) FROM documents").fetchone()[0],
        "users":     db.execute("SELECT COUNT(*) FROM users").fetchone()[0],
        "timestamp": _utcnow().isoformat()
    })

# Auth
@app.route("/api/auth/register",methods=["POST"])
def register():
    d=request.json or {}
    name=(d.get("name") or "").strip(); email=(d.get("email") or "").strip().lower()
    pw=(d.get("password") or "").strip(); dept=(d.get("department") or "Computer Science").strip()
    if not all([name,email,pw]): return jsonify({"error":"All fields required"}),400
    if len(pw)<6: return jsonify({"error":"Password min 6 chars"}),400
    db=get_db()
    if db.execute("SELECT id FROM users WHERE email=?",(email,)).fetchone():
        return jsonify({"error":"Email already registered"}),409
    h=bcrypt.hashpw(pw.encode(),bcrypt.gensalt()).decode()
    cur=db.execute("INSERT INTO users (name,email,password,role,department) VALUES (?,?,?,?,?)",
                   (name,email,h,"faculty",dept))
    db.commit(); uid=cur.lastrowid
    return jsonify({"token":make_token(uid,"faculty"),
                    "user":{"id":uid,"name":name,"email":email,"role":"faculty","department":dept,"bio":""}}),201

@app.route("/api/auth/login",methods=["POST"])
def login():
    d=request.json or {}
    email=(d.get("email") or "").strip().lower(); pw=(d.get("password") or "").strip()
    if not email or not pw: return jsonify({"error":"Email and password required"}),400
    db=get_db()
    u=db.execute("SELECT * FROM users WHERE email=? AND active=1",(email,)).fetchone()
    if not u or not bcrypt.checkpw(pw.encode(),u["password"].encode()):
        return jsonify({"error":"Invalid credentials"}),401
    u=dict(u)
    db.execute("UPDATE users SET last_login=CURRENT_TIMESTAMP WHERE id=?",(u["id"],))
    db.execute("INSERT INTO activity_log (user_id,action) VALUES (?,?)",(u["id"],"login"))
    db.commit()
    return jsonify({"token":make_token(u["id"],u["role"]),
                    "user":{k:u[k] for k in ("id","name","email","role","department","bio")}})

@app.route("/api/auth/me")
@require_auth
def get_me():
    db=get_db()
    u=db.execute("SELECT id,name,email,role,department,bio,created_at,last_login FROM users WHERE id=?",(g.user_id,)).fetchone()
    return jsonify(dict(u)) if u else (jsonify({"error":"Not found"}),404)

@app.route("/api/auth/change-password",methods=["PUT"])
@require_auth
def change_pw():
    d=request.json or {}
    old=(d.get("old_password") or "").strip(); new=(d.get("new_password") or "").strip()
    if not old or not new or len(new)<6: return jsonify({"error":"Invalid input"}),400
    db=get_db()
    u=db.execute("SELECT * FROM users WHERE id=?",(g.user_id,)).fetchone()
    if not bcrypt.checkpw(old.encode(),u["password"].encode()): return jsonify({"error":"Old password wrong"}),401
    db.execute("UPDATE users SET password=? WHERE id=?",(bcrypt.hashpw(new.encode(),bcrypt.gensalt()).decode(),g.user_id))
    db.commit(); return jsonify({"message":"Password changed"})

# Profile
@app.route("/api/profile",methods=["PUT"])
@require_auth
def update_profile():
    d=request.json or {}; db=get_db(); sets,params=[],[]
    for f in ["name","department","bio"]:
        if f in d: sets.append(f"{f}=?"); params.append(d[f])
    if not sets: return jsonify({"error":"Nothing to update"}),400
    params.append(g.user_id)
    db.execute(f"UPDATE users SET {','.join(sets)} WHERE id=?",params); db.commit()
    u=db.execute("SELECT id,name,email,role,department,bio FROM users WHERE id=?",(g.user_id,)).fetchone()
    return jsonify({"message":"Updated","user":dict(u)})

# Documents
def _strip(docs):
    for d in docs: d.pop("text_content",None); d.pop("embedding",None)
    return docs

@app.route("/api/documents")
@require_auth
def list_docs():
    db=get_db(); query=request.args.get("q","").strip()
    page=max(1,int(request.args.get("page",1))); per=20
    filters={k:request.args.get(k,"") for k in ("course","year","type","level","author")}
    if query or any(filters.values()):
        results=search_docs(query,filters,g.user_id); total=len(results)
        items=_strip(results[(page-1)*per:page*per])
    else:
        total=db.execute("SELECT COUNT(*) FROM documents WHERE is_public=1").fetchone()[0]
        rows=db.execute(
            "SELECT d.*,u.name as uploader_name,GROUP_CONCAT(t.tag,',') as tags "
            "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
            "LEFT JOIN tags t ON t.document_id=d.id WHERE d.is_public=1 "
            "GROUP BY d.id ORDER BY d.created_at DESC LIMIT ? OFFSET ?",
            (per,(page-1)*per)).fetchall()
        items=_strip([dict(r) for r in rows])
    return jsonify({"documents":items,"total":total,"page":page,"per_page":per})

@app.route("/api/documents/<int:doc_id>")
@require_auth
def get_doc(doc_id):
    db=get_db()
    doc=db.execute(
        "SELECT d.*,u.name as uploader_name,u.email as uploader_email,"
        "GROUP_CONCAT(t.tag,',') as tags FROM documents d "
        "LEFT JOIN users u ON d.uploader_id=u.id "
        "LEFT JOIN tags t ON t.document_id=d.id WHERE d.id=? GROUP BY d.id",(doc_id,)).fetchone()
    if not doc: return jsonify({"error":"Not found"}),404
    db.execute("UPDATE documents SET view_count=view_count+1 WHERE id=?",(doc_id,))
    db.execute("INSERT INTO activity_log (user_id,action,document_id) VALUES (?,?,?)",(g.user_id,"view",doc_id))
    db.commit()
    bm=db.execute("SELECT id FROM bookmarks WHERE user_id=? AND document_id=?",(g.user_id,doc_id)).fetchone()
    d=dict(doc); d.pop("text_content",None); d.pop("embedding",None)
    d["bookmarked"]=bool(bm); return jsonify(d)

@app.route("/api/documents/upload",methods=["POST"])
@require_auth
def upload_doc():
    if "file" not in request.files: return jsonify({"error":"No file"}),400
    file=request.files["file"]
    if not file.filename: return jsonify({"error":"Empty filename"}),400
    ext=Path(file.filename).suffix.lower()
    if ext not in ALLOWED: return jsonify({"error":f"Type {ext} not allowed"}),400
    safe=f"{uuid.uuid4().hex}{ext}"; dest=UPLOAD_DIR/safe
    file.save(str(dest)); size=dest.stat().st_size
    if size>MAX_BYTES: dest.unlink(missing_ok=True); return jsonify({"error":"File too large"}),413
    title=request.form.get("title",file.filename)
    desc=request.form.get("description","")
    cc=request.form.get("course_code",""); al=request.form.get("academic_level","")
    rt=request.form.get("resource_type",""); ay=request.form.get("academic_year","")
    pub=int(request.form.get("is_public",1))
    text=extract_text(dest,ext)
    if not (cc and al and rt):
        cls=auto_classify(text,file.filename)
        cc=cc or cls["course_code"]; al=al or cls["academic_level"]; rt=rt or cls["resource_type"]
    emb=get_embedding(f"{title} {desc} {text[:3000]}")
    tags=extract_tags(f"{title} {desc} {text}")
    db=get_db()
    cur=db.execute(
        "INSERT INTO documents (title,description,filename,original_name,file_type,file_size,"
        "course_code,academic_level,resource_type,academic_year,uploader_id,text_content,embedding,is_public) "
        "VALUES (?,?,?,?,?,?,?,?,?,?,?,?,?,?)",
        (title,desc,safe,file.filename,ext,size,cc,al,rt,ay,g.user_id,text,
         json.dumps(emb) if emb else "",pub))
    did=cur.lastrowid
    for tag in tags: db.execute("INSERT INTO tags (document_id,tag) VALUES (?,?)",(did,tag))
    db.execute("INSERT INTO activity_log (user_id,action,document_id) VALUES (?,?,?)",(g.user_id,"upload",did))
    db.commit()
    return jsonify({"message":"Uploaded","id":did,"tags":tags,
                    "classification":{"course_code":cc,"academic_level":al,"resource_type":rt}}),201

@app.route("/api/documents/<int:doc_id>",methods=["PUT"])
@require_auth
def update_doc(doc_id):
    db=get_db(); doc=db.execute("SELECT * FROM documents WHERE id=?",(doc_id,)).fetchone()
    if not doc: return jsonify({"error":"Not found"}),404
    if dict(doc)["uploader_id"]!=g.user_id and g.user_role!="admin": return jsonify({"error":"Forbidden"}),403
    d=request.json or {}; sets,params=[],[]
    for f in ["title","description","course_code","academic_level","resource_type","academic_year","is_public"]:
        if f in d: sets.append(f"{f}=?"); params.append(d[f])
    if not sets: return jsonify({"error":"Nothing to update"}),400
    sets.append("updated_at=CURRENT_TIMESTAMP"); params.append(doc_id)
    db.execute(f"UPDATE documents SET {','.join(sets)} WHERE id=?",params); db.commit()
    return jsonify({"message":"Updated"})

@app.route("/api/documents/<int:doc_id>",methods=["DELETE"])
@require_auth
def delete_doc(doc_id):
    db=get_db(); doc=db.execute("SELECT * FROM documents WHERE id=?",(doc_id,)).fetchone()
    if not doc: return jsonify({"error":"Not found"}),404
    d=dict(doc)
    if d["uploader_id"]!=g.user_id and g.user_role!="admin": return jsonify({"error":"Forbidden"}),403
    try: (UPLOAD_DIR/d["filename"]).unlink(missing_ok=True)
    except: pass
    db.execute("DELETE FROM documents WHERE id=?",(doc_id,)); db.commit()
    return jsonify({"message":"Deleted"})

@app.route("/api/documents/<int:doc_id>/download")
@require_auth
def download_doc(doc_id):
    db=get_db(); doc=db.execute("SELECT * FROM documents WHERE id=?",(doc_id,)).fetchone()
    if not doc: return jsonify({"error":"Not found"}),404
    d=dict(doc)
    db.execute("UPDATE documents SET download_count=download_count+1 WHERE id=?",(doc_id,))
    db.execute("INSERT INTO activity_log (user_id,action,document_id) VALUES (?,?,?)",(g.user_id,"download",doc_id))
    db.commit()
    return send_from_directory(str(UPLOAD_DIR),d["filename"],as_attachment=True,download_name=d["original_name"])

@app.route("/api/documents/<int:doc_id>/recommendations")
@require_auth
def recs(doc_id):
    r=get_recs(doc_id)
    for x in r: x.pop("text_content",None); x.pop("embedding",None)
    return jsonify({"recommendations":r})

@app.route("/api/documents/my")
@require_auth
def my_docs():
    db=get_db()
    rows=db.execute(
        "SELECT d.*,u.name as uploader_name,GROUP_CONCAT(t.tag,',') as tags "
        "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
        "LEFT JOIN tags t ON t.document_id=d.id WHERE d.uploader_id=? "
        "GROUP BY d.id ORDER BY d.created_at DESC",(g.user_id,)).fetchall()
    return jsonify({"documents":_strip([dict(r) for r in rows])})

# Bookmarks
@app.route("/api/bookmarks")
@require_auth
def bookmarks():
    db=get_db()
    rows=db.execute(
        "SELECT d.*,u.name as uploader_name,GROUP_CONCAT(t.tag,',') as tags "
        "FROM bookmarks b JOIN documents d ON b.document_id=d.id "
        "LEFT JOIN users u ON d.uploader_id=u.id LEFT JOIN tags t ON t.document_id=d.id "
        "WHERE b.user_id=? GROUP BY d.id ORDER BY b.created_at DESC",(g.user_id,)).fetchall()
    return jsonify({"documents":_strip([dict(r) for r in rows])})

@app.route("/api/bookmarks/<int:doc_id>",methods=["POST"])
@require_auth
def toggle_bm(doc_id):
    db=get_db()
    ex=db.execute("SELECT id FROM bookmarks WHERE user_id=? AND document_id=?",(g.user_id,doc_id)).fetchone()
    if ex: db.execute("DELETE FROM bookmarks WHERE user_id=? AND document_id=?",(g.user_id,doc_id)); db.commit(); return jsonify({"bookmarked":False})
    db.execute("INSERT INTO bookmarks (user_id,document_id) VALUES (?,?)",(g.user_id,doc_id)); db.commit()
    return jsonify({"bookmarked":True})

# Comments
@app.route("/api/documents/<int:doc_id>/comments")
@require_auth
def get_comments(doc_id):
    db=get_db()
    rows=db.execute(
        "SELECT c.*,u.name as user_name FROM comments c "
        "JOIN users u ON c.user_id=u.id WHERE c.document_id=? ORDER BY c.created_at ASC",(doc_id,)).fetchall()
    return jsonify({"comments":[dict(r) for r in rows]})

@app.route("/api/documents/<int:doc_id>/comments",methods=["POST"])
@require_auth
def add_comment(doc_id):
    d=request.json or {}; content=(d.get("content") or "").strip()
    if not content: return jsonify({"error":"Empty comment"}),400
    db=get_db()
    cur=db.execute("INSERT INTO comments (document_id,user_id,content) VALUES (?,?,?)",(doc_id,g.user_id,content))
    db.commit()
    row=db.execute("SELECT c.*,u.name as user_name FROM comments c JOIN users u ON c.user_id=u.id WHERE c.id=?",(cur.lastrowid,)).fetchone()
    return jsonify(dict(row)),201

@app.route("/api/comments/<int:cid>",methods=["DELETE"])
@require_auth
def del_comment(cid):
    db=get_db(); row=db.execute("SELECT * FROM comments WHERE id=?",(cid,)).fetchone()
    if not row: return jsonify({"error":"Not found"}),404
    if dict(row)["user_id"]!=g.user_id and g.user_role!="admin": return jsonify({"error":"Forbidden"}),403
    db.execute("DELETE FROM comments WHERE id=?",(cid,)); db.commit()
    return jsonify({"message":"Deleted"})

# Tags & Faculty
@app.route("/api/tags/popular")
@require_auth
def pop_tags():
    db=get_db()
    rows=db.execute("SELECT tag,COUNT(*) as cnt FROM tags GROUP BY tag ORDER BY cnt DESC LIMIT 40").fetchall()
    return jsonify({"tags":[dict(r) for r in rows]})

@app.route("/api/faculty")
@require_auth
def faculty():
    db=get_db()
    rows=db.execute(
        "SELECT u.id,u.name,u.department,u.bio,COUNT(d.id) as doc_count "
        "FROM users u LEFT JOIN documents d ON d.uploader_id=u.id "
        "WHERE u.active=1 GROUP BY u.id ORDER BY doc_count DESC").fetchall()
    return jsonify({"faculty":[dict(u) for u in rows]})

# Admin
@app.route("/api/admin/stats")
@require_admin
def admin_stats():
    db=get_db()
    return jsonify({
        "total_documents": db.execute("SELECT COUNT(*) FROM documents").fetchone()[0],
        "total_users":     db.execute("SELECT COUNT(*) FROM users").fetchone()[0],
        "total_downloads": db.execute("SELECT COALESCE(SUM(download_count),0) FROM documents").fetchone()[0],
        "total_views":     db.execute("SELECT COALESCE(SUM(view_count),0) FROM documents").fetchone()[0],
        "total_searches":  db.execute("SELECT COUNT(*) FROM search_logs").fetchone()[0],
        "total_bookmarks": db.execute("SELECT COUNT(*) FROM bookmarks").fetchone()[0],
        "total_comments":  db.execute("SELECT COUNT(*) FROM comments").fetchone()[0],
        "by_type":    [dict(r) for r in db.execute("SELECT resource_type,COUNT(*) cnt FROM documents GROUP BY resource_type ORDER BY cnt DESC").fetchall()],
        "by_course":  [dict(r) for r in db.execute("SELECT course_code,COUNT(*) cnt FROM documents WHERE course_code!='' GROUP BY course_code ORDER BY cnt DESC LIMIT 10").fetchall()],
        "by_level":   [dict(r) for r in db.execute("SELECT academic_level,COUNT(*) cnt FROM documents WHERE academic_level!='' GROUP BY academic_level ORDER BY cnt DESC").fetchall()],
        "recent_uploads": [dict(r) for r in db.execute("SELECT d.id,d.title,u.name as uploader_name,d.created_at,d.file_type,d.download_count FROM documents d JOIN users u ON d.uploader_id=u.id ORDER BY d.created_at DESC LIMIT 10").fetchall()],
        "daily_uploads":  [dict(r) for r in db.execute("SELECT DATE(created_at) date,COUNT(*) cnt FROM documents GROUP BY DATE(created_at) ORDER BY date DESC LIMIT 30").fetchall()],
        "daily_searches": [dict(r) for r in db.execute("SELECT DATE(created_at) date,COUNT(*) cnt FROM search_logs GROUP BY DATE(created_at) ORDER BY date DESC LIMIT 30").fetchall()],
    })

@app.route("/api/admin/users")
@require_admin
def admin_users():
    db=get_db()
    rows=db.execute(
        "SELECT u.id,u.name,u.email,u.role,u.department,u.active,u.created_at,u.last_login,"
        "COUNT(d.id) as doc_count FROM users u LEFT JOIN documents d ON d.uploader_id=u.id "
        "GROUP BY u.id ORDER BY u.created_at DESC").fetchall()
    return jsonify({"users":[dict(u) for u in rows]})

@app.route("/api/admin/users/<int:uid>",methods=["PUT"])
@require_admin
def admin_upd_user(uid):
    d=request.json or {}; db=get_db(); sets,params=[],[]
    for f in ["role","active","department","name"]:
        if f in d: sets.append(f"{f}=?"); params.append(d[f])
    if not sets: return jsonify({"error":"Nothing to update"}),400
    params.append(uid); db.execute(f"UPDATE users SET {','.join(sets)} WHERE id=?",params); db.commit()
    return jsonify({"message":"Updated"})

@app.route("/api/admin/users/<int:uid>",methods=["DELETE"])
@require_admin
def admin_del_user(uid):
    if uid==g.user_id: return jsonify({"error":"Cannot deactivate yourself"}),400
    db=get_db(); db.execute("UPDATE users SET active=0 WHERE id=?",(uid,)); db.commit()
    return jsonify({"message":"Deactivated"})

@app.route("/api/admin/documents")
@require_admin
def admin_all_docs():
    db=get_db()
    rows=db.execute(
        "SELECT d.*,u.name as uploader_name,GROUP_CONCAT(t.tag,',') as tags "
        "FROM documents d LEFT JOIN users u ON d.uploader_id=u.id "
        "LEFT JOIN tags t ON t.document_id=d.id GROUP BY d.id ORDER BY d.created_at DESC").fetchall()
    return jsonify({"documents":_strip([dict(r) for r in rows])})

@app.route("/api/admin/reports/usage")
@require_admin
def usage_report():
    db=get_db()
    return jsonify({
        "top_documents": [dict(r) for r in db.execute("SELECT id,title,download_count,view_count FROM documents ORDER BY download_count DESC LIMIT 10").fetchall()],
        "top_users":     [dict(r) for r in db.execute("SELECT u.id,u.name,u.email,COUNT(d.id) as doc_count FROM users u LEFT JOIN documents d ON d.uploader_id=u.id GROUP BY u.id ORDER BY doc_count DESC LIMIT 10").fetchall()],
        "daily_searches":[dict(r) for r in db.execute("SELECT DATE(created_at) date,COUNT(*) searches FROM search_logs GROUP BY DATE(created_at) ORDER BY date DESC LIMIT 30").fetchall()],
    })

@app.route("/api/admin/activity")
@require_admin
def activity():
    db=get_db()
    rows=db.execute(
        "SELECT a.*,u.name as user_name,d.title as doc_title FROM activity_log a "
        "LEFT JOIN users u ON a.user_id=u.id LEFT JOIN documents d ON a.document_id=d.id "
        "ORDER BY a.created_at DESC LIMIT 100").fetchall()
    return jsonify({"activity":[dict(r) for r in rows]})


if __name__ == "__main__":
    init_db()
    log.info(f"SLU Repo starting on port {PORT}")
    log.info(f"HF Space: {_HF_URL or 'NOT CONFIGURED — using TF-IDF fallback'}")
    app.run(host="0.0.0.0", port=PORT, debug=False)
